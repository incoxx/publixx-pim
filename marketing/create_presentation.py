#!/usr/bin/env python3
"""
anyPIM Agency Presentation Generator
Generates a professional Dark Mode PowerPoint presentation for agency partners.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Design Constants ───────────────────────────────────────────────────────
BG_COLOR = RGBColor(0x1a, 0x1a, 0x2e)
BG_CARD = RGBColor(0x22, 0x22, 0x3a)
ACCENT = RGBColor(0x00, 0xd4, 0xff)
ACCENT_ORANGE = RGBColor(0xff, 0x6b, 0x35)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xb0, 0xb0, 0xb0)
DARK_GRAY = RGBColor(0x2a, 0x2a, 0x44)
GREEN = RGBColor(0x00, 0xe6, 0x76)
RED = RGBColor(0xff, 0x44, 0x44)
YELLOW = RGBColor(0xff, 0xd7, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"


def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                       line_spacing=1.3):
    """lines: list of (text, color, bold) or just str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, clr, bld = line, color, bold
        else:
            txt, clr, bld = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = FONT_BODY
        p.alignment = alignment
        p.line_spacing = Pt(font_size * line_spacing)
        p.space_after = Pt(2)
    return txBox


def add_card(slide, left, top, width, height, border_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def add_icon_card(slide, left, top, width, height, icon, title, desc,
                  border_color=ACCENT, icon_size=32, title_size=16, desc_size=12):
    add_card(slide, left, top, width, height, border_color)
    add_text_box(slide, left, top + Inches(0.15), width, Inches(0.5),
                 icon, font_size=icon_size, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.65), width - Inches(0.2), Inches(0.35),
                 title, font_size=title_size, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(1.0), width - Inches(0.2), height - Inches(1.1),
                 desc, font_size=desc_size, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def add_section_header(slide, title, subtitle=None, icon=None):
    set_slide_bg(slide)
    # Top accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    full_title = f"{icon}  {title}" if icon else title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 full_title, font_size=36, color=ACCENT, bold=True, font_name=FONT_TITLE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.45),
                     subtitle, font_size=18, color=LIGHT_GRAY)


def add_bullet_card(slide, left, top, width, height, title, bullets, border_color=ACCENT):
    add_card(slide, left, top, width, height, border_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=ACCENT, bold=True)
    lines = []
    for b in bullets:
        if isinstance(b, tuple):
            lines.append((f"  {b[0]}", b[1], False))
        else:
            lines.append((f"  {b}", LIGHT_GRAY, False))
    add_multiline_text(slide, left + Inches(0.2), top + Inches(0.5),
                       width - Inches(0.4), height - Inches(0.6),
                       lines, font_size=13, line_spacing=1.35)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Titelfolie"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Decorative top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # anyPIM logo text
    add_text_box(slide, Inches(0), Inches(1.8), SLIDE_WIDTH, Inches(1.2),
                 "anyPIM", font_size=72, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Tagline
    add_text_box(slide, Inches(0), Inches(3.0), SLIDE_WIDTH, Inches(0.6),
                 "Das moderne Product Information Management", font_size=28,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(0), Inches(3.7), SLIDE_WIDTH, Inches(0.5),
                 "Partnerpraesentation fuer Agenturen", font_size=20,
                 color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(4.5), Inches(4.5), Inches(4.333), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # Footer
    add_text_box(slide, Inches(0), Inches(5.3), SLIDE_WIDTH, Inches(0.4),
                 "incoxx GmbH  |  Software & Support Partner", font_size=14,
                 color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def slide_02_what_is(prs):
    """Was ist anyPIM?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Was ist anyPIM?", "Modernes PIM -- sofort einsatzbereit, kein Framework")

    # Left: Description
    add_multiline_text(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.5), [
        ("anyPIM ist ein vollstaendiges Product Information", WHITE, False),
        ("Management System -- kein Framework, das erst", WHITE, False),
        ("zusammengebaut werden muss.", WHITE, False),
        ("", WHITE, False),
        ("Entwickelt fuer Teams, die sofort produktiv sein", WHITE, False),
        ("wollen, ohne ein PHP-Entwicklerteam zu brauchen.", WHITE, False),
    ], font_size=17, line_spacing=1.4)

    # Right: Key facts cards
    facts = [
        ("Backend", "Laravel 11 + PHP 8.3", ACCENT),
        ("Frontend", "Vue 3 + Inertia.js", ACCENT),
        ("API", "285+ REST Endpoints", GREEN),
        ("Architektur", "EAV -- Schema per UI", ACCENT),
        ("Suche", "Elasticsearch + PQL", ACCENT),
        ("Setup", "Minuten, nicht Monate", GREEN),
    ]
    for i, (label, value, color) in enumerate(facts):
        row = i // 2
        col = i % 2
        x = Inches(7.2) + col * Inches(2.9)
        y = Inches(1.7) + row * Inches(1.6)
        add_card(slide, x, y, Inches(2.7), Inches(1.35), color)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.4), Inches(0.35),
                     label, font_size=13, color=LIGHT_GRAY, bold=False)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.4), Inches(0.6),
                     value, font_size=16, color=WHITE, bold=True)


def slide_03_features_overview(prs):
    """Kernfunktionen Uebersicht"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Kernfunktionen", "Alles was ein modernes PIM braucht -- und mehr")

    features = [
        ("*", "Produkte", "Datenmodelle, Varianten,\nVererbung, Versionierung"),
        ("+", "Attribute", "45+ Feldtypen, EAV,\nValidierung, Gruppen"),
        ("#", "Hierarchien", "Master + Output Trees,\nMulti-Channel Ready"),
        ("~", "Media", "DAM, Bildbearbeitung,\nFormate, Metadaten"),
        ("$", "Preise", "Staffelpreise, Waehrungen,\nKundengruppen"),
        ("&", "Relationen", "Cross-Sell, Up-Sell,\nBundles, Zubehoer"),
        ("?", "Suche & PQL", "Fuzzy, Phonetisch,\nQuery Language"),
        ("!", "Workflows", "Status, Transitions,\nBenachrichtigungen"),
    ]

    for i, (icon, title, desc) in enumerate(features):
        col = i % 4
        row = i // 4
        x = Inches(0.6) + col * Inches(3.1)
        y = Inches(1.7) + row * Inches(2.6)
        add_icon_card(slide, x, y, Inches(2.8), Inches(2.3),
                      icon, title, desc, ACCENT, icon_size=28, title_size=15, desc_size=12)


def slide_04_pdf_designer(prs):
    """PDF Template Designer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "PDF Template Designer",
                       "Alleinstellungsmerkmal -- Visueller PDF-Designer direkt im PIM")

    # USP badge
    add_text_box(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.4),
                 "USP", font_size=16, color=ACCENT_ORANGE, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # Left column: Features
    add_bullet_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8),
                    "Features", [
                        "Visueller Drag & Drop Template-Editor",
                        "Platzhalter: Produktname, Attribute, Preise, Bilder",
                        "Loop-Bloecke fuer Listen und Tabellen",
                        "Formate: A4, A3, Letter, Custom",
                        "Output: PDF und DOCX",
                        "Bulk-Rendering: Hunderte PDFs auf Knopfdruck",
                        "Vorschau direkt im Editor",
                        "Eigene Schriften und Logos einbindbar",
                    ])

    # Right column: Comparison
    add_card(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(2.2), RED)
    add_text_box(slide, Inches(7.2), Inches(1.85), Inches(5.0), Inches(0.35),
                 "Pimcore: KEIN PDF Designer", font_size=16, color=RED, bold=True)
    add_multiline_text(slide, Inches(7.2), Inches(2.3), Inches(5.0), Inches(1.4), [
        ("Erfordert Drittanbieter:", LIGHT_GRAY, False),
        ("  PDFreactor (kommerziell, ab 2.950 EUR)", WHITE, False),
        ("  Oder: Manuelle HTML/CSS/PHP Templates", WHITE, False),
        ("  Kein visueller Editor verfuegbar", WHITE, False),
    ], font_size=14, line_spacing=1.4)

    # Use case
    add_card(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(2.3), GREEN)
    add_text_box(slide, Inches(7.2), Inches(4.35), Inches(5.0), Inches(0.35),
                 "Anwendungsbeispiele", font_size=16, color=GREEN, bold=True)
    add_multiline_text(slide, Inches(7.2), Inches(4.8), Inches(5.0), Inches(1.5), [
        ("  Produktdatenblaetter", WHITE, False),
        ("  Preislisten mit Staffelpreisen", WHITE, False),
        ("  Technische Spezifikationen", WHITE, False),
        ("  Katalogseiten fuer Print", WHITE, False),
    ], font_size=14, line_spacing=1.4)


def slide_05_indesign(prs):
    """InDesign Export"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "InDesign Export",
                       "Alleinstellungsmerkmal -- Komplettes Layout-Paket als ZIP")

    add_text_box(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.4),
                 "USP", font_size=16, color=ACCENT_ORANGE, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # Left: How it works
    add_bullet_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8),
                    "So funktioniert es", [
                        "Template-basierte InDesign-Pakete",
                        "ZIP mit JSX-Script + JSON-Daten + Assets",
                        "Automatische Seitenerstellung in InDesign",
                        "Tabellen mit vollstaendigem Styling",
                        "Shapes, Textrahmen, Bilder",
                        "Fonts und Farben aus dem Template",
                        "Bulk-Export fuer grosse Kataloge",
                        "Nahtlose Integration in Print-Workflows",
                    ])

    # Right: Pimcore comparison
    add_card(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.8), RED)
    add_text_box(slide, Inches(7.2), Inches(1.85), Inches(5.0), Inches(0.35),
                 "Pimcore: Nur ueber teure Drittanbieter", font_size=16, color=RED, bold=True)
    add_multiline_text(slide, Inches(7.2), Inches(2.4), Inches(5.0), Inches(4.0), [
        ("Verfuegbare Optionen:", LIGHT_GRAY, True),
        ("", WHITE, False),
        ("PimPrint", WHITE, True),
        ("  Lizenzkosten + Implementierung", LIGHT_GRAY, False),
        ("  Eigenstaendiges Plugin", LIGHT_GRAY, False),
        ("", WHITE, False),
        ("Pagination.com", WHITE, True),
        ("  Externer SaaS-Dienst", LIGHT_GRAY, False),
        ("  Zusaetzliche laufende Kosten", LIGHT_GRAY, False),
        ("", WHITE, False),
        ("EasyCatalog / Factory.dev", WHITE, True),
        ("  Weitere kommerzielle Plugins", LIGHT_GRAY, False),
        ("", WHITE, False),
        ("Keine dieser Loesungen ist in", RED, True),
        ("Pimcore enthalten!", RED, True),
    ], font_size=14, line_spacing=1.3)


def slide_06_report_designer(prs):
    """Report Designer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Report Designer",
                       "Alleinstellungsmerkmal -- Berichte erstellen ohne Entwickler")

    add_text_box(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.4),
                 "USP", font_size=16, color=ACCENT_ORANGE, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # Features in 3 cards
    cards = [
        ("Visualisierungen", [
            "Tabellen mit Sortierung & Filter",
            "Diagramme (Bar, Line, Pie)",
            "Summary-Karten mit KPIs",
            "Drag & Drop Feldauswahl",
        ], ACCENT),
        ("Export & Automation", [
            "PDF, DOCX, CSV, XLSX Export",
            "Cron-basiertes Scheduling",
            "Automatischer E-Mail-Versand",
            "SearchProfile-Integration",
        ], GREEN),
        ("vs. Pimcore", [
            "Pimcore: 'Statistics Explorer'",
            "  Erst seit Version 2025.1",
            "  Nur einfache Reports",
            "  Kein Scheduling/Export",
        ], RED),
    ]

    for i, (title, bullets, color) in enumerate(cards):
        x = Inches(0.8) + i * Inches(4.1)
        add_bullet_card(slide, x, Inches(1.7), Inches(3.8), Inches(4.0), title, bullets, color)


def slide_07_project_management(prs):
    """Projekt- & Workflow-Management"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Projekt- & Workflow-Management",
                       "Alleinstellungsmerkmal -- Integriertes Projektmanagement im PIM")

    add_text_box(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.4),
                 "USP", font_size=16, color=ACCENT_ORANGE, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # Left: Project features
    add_bullet_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.2),
                    "Projekte", [
                        "Projekte mit Teams & Mitgliedern",
                        "Aufgaben: open -> in_progress -> review -> done",
                        "Prioritaeten & Deadlines",
                        "Bulk-Tasks fuer Massenbearbeitung",
                    ], ACCENT)

    add_bullet_card(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.3),
                    "Workflows", [
                        "Produkt-Status-Workflows (Draft -> Review -> Published)",
                        "Automatische Benachrichtigungen",
                        "Berechtigungsbasierte Transitions",
                        "Audit Trail & History",
                    ], ACCENT)

    # Right: vs Pimcore
    add_card(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.8), RED)
    add_text_box(slide, Inches(7.2), Inches(1.85), Inches(5.0), Inches(0.4),
                 "Pimcore: KEIN Projektmanagement", font_size=18, color=RED, bold=True)
    add_multiline_text(slide, Inches(7.2), Inches(2.5), Inches(5.0), Inches(3.5), [
        ("Pimcore hat einen Workflow-Engine 2.0,", LIGHT_GRAY, False),
        ("aber:", WHITE, False),
        ("", WHITE, False),
        ("  Keine Projekte", WHITE, False),
        ("  Keine Aufgabenverwaltung", WHITE, False),
        ("  Keine Team-Zuordnung", WHITE, False),
        ("  Keine Deadlines/Prioritaeten", WHITE, False),
        ("  Keine Bulk-Task-Erstellung", WHITE, False),
        ("", WHITE, False),
        ("Agenturen muessen auf externe Tools", LIGHT_GRAY, False),
        ("wie Jira, Asana oder Trello ausweichen.", LIGHT_GRAY, False),
        ("", WHITE, False),
        ("anyPIM: Alles in einem System.", GREEN, True),
    ], font_size=15, line_spacing=1.3)


def slide_08_pql(prs):
    """PQL -- Product Query Language"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "PQL -- Product Query Language",
                       "Alleinstellungsmerkmal -- Einzigartig in der PIM-Welt")

    add_text_box(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.4),
                 "USP", font_size=16, color=ACCENT_ORANGE, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # PQL Examples (code-like)
    examples = [
        'name FUZZY "Shraubbe"  -- findet "Schraube"',
        'name SOUNDS_LIKE "Mayer"  -- findet Meyer, Meier, Maier',
        'SEARCH_FIELDS(name^3, description) "Motor"',
        'price > 10 AND category = "Elektronik"',
        'status IN ("active", "draft") ORDER BY relevance',
    ]

    add_card(slide, Inches(0.8), Inches(1.7), Inches(7.0), Inches(4.8), ACCENT)
    add_text_box(slide, Inches(1.0), Inches(1.85), Inches(6.5), Inches(0.35),
                 "PQL Beispiele", font_size=16, color=ACCENT, bold=True)

    lines = []
    for ex in examples:
        lines.append((f"  {ex}", WHITE, False))
        lines.append(("", WHITE, False))
    add_multiline_text(slide, Inches(1.0), Inches(2.3), Inches(6.5), Inches(4.0),
                       lines, font_size=14, line_spacing=1.3)

    # Right: Features
    add_bullet_card(slide, Inches(8.3), Inches(1.7), Inches(4.5), Inches(4.8),
                    "Funktionen", [
                        "SQL-aehnliche Syntax",
                        "FUZZY: Tippfehler-tolerant",
                        "SOUNDS_LIKE: Phonetische Suche",
                        "  Koelner Phonetik (DE)",
                        "  Soundex (EN)",
                        "Gewichtete Feldsuche",
                        "Relevanz-Ranking",
                        "Komplex verschachtelbar",
                        "",
                        ("Pimcore: Hat NICHTS", RED, True),
                        ("Vergleichbares!", RED, True),
                    ], GREEN)


def slide_09_api_designer(prs):
    """API Designer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "API Designer",
                       "Alleinstellungsmerkmal -- Visuelle API-Konfiguration")

    add_text_box(slide, Inches(10.5), Inches(0.4), Inches(2.5), Inches(0.4),
                 "USP", font_size=16, color=ACCENT_ORANGE, bold=True,
                 alignment=PP_ALIGN.RIGHT)

    # Left: Features
    add_bullet_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8),
                    "Visueller Endpoint-Builder", [
                        "API-Endpoints per Drag & Drop konfigurieren",
                        "Felder, Filter und Sortierung festlegen",
                        "API-Key-Management mit Berechtigungen",
                        "Rate Limiting pro Key/Endpoint",
                        "Streaming fuer grosse Datenmengen",
                        "Import- und Export-Endpoints",
                        "Webhook-Integration",
                        "Automatische API-Dokumentation",
                    ], ACCENT)

    # Right: vs Pimcore
    add_card(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(2.2), RED)
    add_text_box(slide, Inches(7.2), Inches(1.85), Inches(5.0), Inches(0.35),
                 "Pimcore: Nur fuer Entwickler", font_size=16, color=RED, bold=True)
    add_multiline_text(slide, Inches(7.2), Inches(2.35), Inches(5.0), Inches(1.5), [
        ("  Data Hub: GraphQL/REST moeglich", LIGHT_GRAY, False),
        ("  Aber: Konfiguration nur per Code", LIGHT_GRAY, False),
        ("  Erfordert Symfony/PHP-Entwickler", LIGHT_GRAY, False),
        ("  Kein visueller Designer", LIGHT_GRAY, False),
    ], font_size=14, line_spacing=1.4)

    # Stats
    add_card(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(2.3), GREEN)
    add_text_box(slide, Inches(7.2), Inches(4.35), Inches(5.0), Inches(0.35),
                 "anyPIM API in Zahlen", font_size=16, color=GREEN, bold=True)
    add_multiline_text(slide, Inches(7.2), Inches(4.8), Inches(5.0), Inches(1.5), [
        ("  285+ REST Endpoints", WHITE, True),
        ("  Vollstaendige OpenAPI-Dokumentation", WHITE, False),
        ("  Import, Export, CRUD, Bulk, Search", WHITE, False),
        ("  Streaming fuer 100K+ Produkte", WHITE, False),
    ], font_size=14, line_spacing=1.4)


def slide_10_bulk_import(prs):
    """Bulk Editor & Import"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Bulk Editor & Import",
                       "Massenbearbeitung und intelligenter Excel-Import")

    # Left: Bulk Editor
    add_bullet_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8),
                    "Spreadsheet-Style Bulk Editor", [
                        "Inline-Bearbeitung wie in Excel",
                        "Bis zu 5.000 Werte pro Transaktion",
                        "Multi-Select fuer Massenoperationen",
                        "Spalten ein-/ausblenden",
                        "Filter und Sortierung",
                        "Undo/Redo Unterstuetzung",
                        "Copy & Paste aus Excel",
                    ], ACCENT)

    # Right: Import
    add_bullet_card(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.8),
                    "14-Tab Excel Import", [
                        "Produkte, Attribute, Preise, Medien ...",
                        "Automatisches Fuzzy Column Matching",
                        "  Erkennt 'Artikelnr.' als 'SKU'",
                        "  Erkennt 'Bezeichnung' als 'Name'",
                        "Validierung vor dem Import",
                        "Fehlerprotokoll mit Zeilenangabe",
                        "Update-Modus: Merge oder Replace",
                        "",
                        ("Pimcore: Nur einfacher CSV Import", RED, False),
                    ], GREEN)


def slide_11_catalog(prs):
    """Katalog & Zugangslinks"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Oeffentlicher Katalog & Zugangslinks",
                       "Produktdaten teilen -- ohne Login-Zwang")

    cards = [
        ("Oeffentlicher Katalog", [
            "Produkte ohne Login teilen",
            "Konfigurierbare Ansichten",
            "Produktvergleich",
            "Merkliste / Favoriten",
            "PDF & Excel Export",
        ], ACCENT),
        ("Zugangslinks", [
            "Temporaere Links fuer Partner",
            "Zeitlich begrenzt",
            "Passwort-geschuetzt moeglich",
            "Tracking: Wer hat was gesehen?",
            "Sofort widerrufbar",
        ], GREEN),
        ("vs. Pimcore", [
            "Portal Engine nur in",
            "  Enterprise Edition (ab 2.000 EUR/M)",
            "Erhebliche Konfiguration noetig",
            "Kein einfacher 'Katalog teilen'-Button",
            "Keine Zugangslinks",
        ], RED),
    ]

    for i, (title, bullets, color) in enumerate(cards):
        x = Inches(0.8) + i * Inches(4.1)
        add_bullet_card(slide, x, Inches(1.7), Inches(3.8), Inches(4.5), title, bullets, color)


def slide_12_i18n(prs):
    """Mehrsprachigkeit & TMS"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Mehrsprachigkeit & Translation Management",
                       "Globale Produkte, lokale Inhalte")

    add_bullet_card(slide, Inches(0.8), Inches(1.7), Inches(3.8), Inches(4.5),
                    "3-Ebenen Internationalisierung", [
                        "UI-Sprache (Backend)",
                        "Produktdaten-Sprache",
                        "Katalog-/Output-Sprache",
                        "Unabhaengig konfigurierbar",
                        "Fallback-Ketten",
                    ], ACCENT)

    add_bullet_card(slide, Inches(5.0), Inches(1.7), Inches(3.8), Inches(4.5),
                    "Translation Workflow", [
                        "XLIFF Export fuer externe TMS",
                        "XLIFF Import mit Merge",
                        "Uebersetzungsstatistiken",
                        "  % pro Sprache und Feld",
                        "Fehlende Uebersetzungen markiert",
                    ], GREEN)

    add_bullet_card(slide, Inches(9.2), Inches(1.7), Inches(3.8), Inches(4.5),
                    "vs. Pimcore", [
                        "Basis: CSV Import/Export",
                        "DeepL-Integration verfuegbar",
                        "Aber: Kein TMS-Connector",
                        "Kein XLIFF Support",
                        "Keine Uebersetzungsstatistiken",
                    ], RED)


def slide_13_comparison(prs):
    """anyPIM vs. Pimcore -- Vergleichstabelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "anyPIM vs. Pimcore",
                       "Feature-Vergleich auf einen Blick")

    rows = [
        ("Setup-Zeit", "Minuten", "Wochen bis Monate"),
        ("PDF Template Designer", "Built-in, visuell", "Nicht vorhanden"),
        ("InDesign Export", "Built-in als ZIP", "Nur Drittanbieter ($$$)"),
        ("Report Designer", "Built-in, Scheduling", "Seit 2025.1, sehr limitiert"),
        ("Projektmanagement", "Projekte + Aufgaben", "Nicht vorhanden"),
        ("PQL Query Language", "Fuzzy + Phonetisch", "Nicht vorhanden"),
        ("API Designer", "Visueller Builder", "Nur per Code"),
        ("Schema-Aenderungen", "UI/API, kein Deploy", "Entwickler noetig"),
        ("Oeffentlicher Katalog", "Built-in, sofort", "Nur Enterprise (2K/M)"),
        ("Bulk Excel Import", "14-Tab, Fuzzy Match", "Einfacher CSV Import"),
        ("Entwickler noetig?", "Nein", "PHP/Symfony-Team"),
        ("Implementation", "Transparent", "Ab 50.000 EUR"),
    ]

    # Table header
    col_x = [Inches(0.8), Inches(5.8), Inches(9.2)]
    col_w = [Inches(4.8), Inches(3.2), Inches(3.8)]
    header_y = Inches(1.65)
    row_h = Inches(0.38)

    # Header row
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.8), header_y, Inches(12.2), row_h)
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = ACCENT
    header_bar.line.fill.background()

    headers = ["Feature", "anyPIM", "Pimcore"]
    for j, h in enumerate(headers):
        add_text_box(slide, col_x[j] + Inches(0.1), header_y + Inches(0.05),
                     col_w[j], Inches(0.3), h, font_size=14, color=BG_COLOR, bold=True)

    # Data rows
    for i, (feature, anypim, pimcore) in enumerate(rows):
        y = header_y + row_h + i * row_h
        bg = BG_CARD if i % 2 == 0 else DARK_GRAY
        row_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0.8), y, Inches(12.2), row_h)
        row_bar.fill.solid()
        row_bar.fill.fore_color.rgb = bg
        row_bar.line.fill.background()

        add_text_box(slide, col_x[0] + Inches(0.1), y + Inches(0.05),
                     col_w[0], Inches(0.28), feature, font_size=12, color=WHITE, bold=False)
        add_text_box(slide, col_x[1] + Inches(0.1), y + Inches(0.05),
                     col_w[1], Inches(0.28), anypim, font_size=12, color=GREEN, bold=True)
        add_text_box(slide, col_x[2] + Inches(0.1), y + Inches(0.05),
                     col_w[2], Inches(0.28), pimcore, font_size=12, color=RED, bold=False)


def slide_14_partner_model(prs):
    """Partnermodell"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Partnermodell",
                       "So funktioniert die Zusammenarbeit")

    # Agency card
    add_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4.8), ACCENT)
    add_text_box(slide, Inches(1.0), Inches(1.85), Inches(5.0), Inches(0.5),
                 "Agentur", font_size=22, color=ACCENT, bold=True)
    add_multiline_text(slide, Inches(1.0), Inches(2.5), Inches(5.0), Inches(3.5), [
        ("Ihre Rolle:", WHITE, True),
        ("", WHITE, False),
        ("  Kundenberatung & Anforderungsanalyse", WHITE, False),
        ("  Datenmodell-Design & Konfiguration", WHITE, False),
        ("  Import & Migration bestehender Daten", WHITE, False),
        ("  Schulung der Endanwender", WHITE, False),
        ("  Template-Erstellung (PDF, InDesign)", WHITE, False),
        ("  Laufende Betreuung & Optimierung", WHITE, False),
        ("", WHITE, False),
        ("  Umsatz durch Services & Beratung", GREEN, True),
    ], font_size=15, line_spacing=1.3)

    # incoxx card
    add_card(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.8), ACCENT_ORANGE)
    add_text_box(slide, Inches(7.2), Inches(1.85), Inches(5.0), Inches(0.5),
                 "incoxx GmbH", font_size=22, color=ACCENT_ORANGE, bold=True)
    add_multiline_text(slide, Inches(7.2), Inches(2.5), Inches(5.0), Inches(3.5), [
        ("Unser Part:", WHITE, True),
        ("", WHITE, False),
        ("  Software-Entwicklung & Updates", WHITE, False),
        ("  Technischer Support (2nd/3rd Level)", WHITE, False),
        ("  Hosting & Infrastruktur (optional)", WHITE, False),
        ("  Onboarding fuer Agentur-Teams", WHITE, False),
        ("  Gemeinsame Kundenakquise", WHITE, False),
        ("  Feature Requests & Roadmap", WHITE, False),
        ("", WHITE, False),
        ("  Partnerpreise & Margen", GREEN, True),
    ], font_size=15, line_spacing=1.3)


def slide_15_next_steps(prs):
    """Naechste Schritte & Kontakt"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top accent line
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text_box(slide, Inches(0), Inches(1.3), SLIDE_WIDTH, Inches(0.7),
                 "Naechste Schritte", font_size=40, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Steps
    steps = [
        ("1", "Live-Demo", "Wir zeigen Ihnen anyPIM im Detail"),
        ("2", "Testzugang", "Ihre eigene Instanz zum Ausprobieren"),
        ("3", "Partnervertrag", "Konditionen und Zusammenarbeit definieren"),
        ("4", "Erstes Kundenprojekt", "Gemeinsam zum Erfolg"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(1.2) + i * Inches(3.0)
        y = Inches(2.5)
        add_card(slide, x, y, Inches(2.7), Inches(1.8), ACCENT if i < 2 else ACCENT_ORANGE)
        add_text_box(slide, x, y + Inches(0.15), Inches(2.7), Inches(0.45),
                     num, font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.6), Inches(2.5), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), y + Inches(1.0), Inches(2.5), Inches(0.6),
                     desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Contact
    add_text_box(slide, Inches(0), Inches(5.0), SLIDE_WIDTH, Inches(0.5),
                 "incoxx GmbH", font_size=24, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(0.4),
                 "www.incoxx.com  |  info@incoxx.com", font_size=16,
                 color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(5.95), SLIDE_WIDTH, Inches(0.4),
                 "Wir freuen uns auf die Zusammenarbeit!", font_size=18,
                 color=WHITE, alignment=PP_ALIGN.CENTER, bold=True)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Generate all slides
    slide_01_title(prs)
    slide_02_what_is(prs)
    slide_03_features_overview(prs)
    slide_04_pdf_designer(prs)
    slide_05_indesign(prs)
    slide_06_report_designer(prs)
    slide_07_project_management(prs)
    slide_08_pql(prs)
    slide_09_api_designer(prs)
    slide_10_bulk_import(prs)
    slide_11_catalog(prs)
    slide_12_i18n(prs)
    slide_13_comparison(prs)
    slide_14_partner_model(prs)
    slide_15_next_steps(prs)

    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "anyPIM_Agentur_Praesentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
